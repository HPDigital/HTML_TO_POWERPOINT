"""
HTML_TO_POWERPOINT
"""

#!/usr/bin/env python
# coding: utf-8

# In[5]:


get_ipython().system('pip install python-pptx')


# In[7]:


from pptx import Presentation
from pptx.util import Inches, Pt
from bs4 import BeautifulSoup

# Cargar el archivo HTML
html_file_path = r"C:\Users\HP\Downloads\Presentacion_html.html"
with open(html_file_path, 'r', encoding='utf-8') as file:
    html_content = file.read()

# Parsear el HTML con BeautifulSoup
soup = BeautifulSoup(html_content, 'html.parser')

# Crear una presentación PowerPoint
prs = Presentation()

# Establecer el tamaño de fuente predeterminado
font_size_h1 = Pt(24)
font_size_h2 = Pt(20)
font_size_h3 = Pt(18)
font_size_p = Pt(14)

# Función para añadir una diapositiva con un título y contenido
def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Título y contenido
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = content

# Función para procesar los elementos HTML y añadirlos a la presentación
def process_html(soup):
    title = ""
    content = ""
    for element in soup.body:
        if element.name == 'h1':
            if title and content:
                add_slide(title, content)
            title = element.text
            content = ""
        elif element.name == 'h2':
            if content:
                content += '\n\n'
            content += element.text
        elif element.name == 'h3':
            if content:
                content += '\n'
            content += element.text
        elif element.name == 'p':
            if content:
                content += '\n'
            content += element.text
        elif element.name == 'ul':
            for li in element.find_all('li', recursive=False):
                if title and content:
                    add_slide(title, content)
                title = li.strong.text if li.strong else "Diapositiva"
                content = "\n".join([li.text for li in li.find_all('li', recursive=False)])

    # Añadir la última diapositiva si queda contenido sin añadir
    if title and content:
        add_slide(title, content)

# Procesar el HTML
process_html(soup)

# Guardar la presentación en un archivo
output_file = r"C:\Users\HP\Downloads\Presentacion_html.pptx"
prs.save(output_file)
print(f"Presentación guardada en {output_file}")


# In[ ]:






if __name__ == "__main__":
    pass
